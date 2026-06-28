# -*- coding: utf-8 -*-
"""
"AI와 일하는 법 — 왜 AI는 강의를 할 수 없는가" 마스터클래스 슬라이드 생성기.
실행: python build_slides.py  ->  AI와_일하는_법.pptx

컨셉(하이브리드):
  - 메타 후크: "이 강의는 필요 없습니다 / 이 슬라이드도 AI가 만들었습니다"
  - 본문 뼈대: AI 사장님 되기 (비개발자 셀러 친화)
  - 마무리 명분: 운전면허 비유 (도구는 바뀌어도 판단은 남는다)
"""

import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- 테마 ----
FONT = "맑은 고딕"
NAVY = RGBColor(0x14, 0x1B, 0x2E)      # 진한 남색 (제목/배경 강조)
INK = RGBColor(0x22, 0x28, 0x33)       # 본문 글자
CORAL = RGBColor(0xFF, 0x5A, 0x4D)     # 강조 (X / 경고 / 포인트)
TEAL = RGBColor(0x16, 0xA3, 0x85)      # 긍정 (O / 추천)
GOLD = RGBColor(0xF5, 0xB3, 0x41)      # 보조 강조
GRAY = RGBColor(0x6B, 0x72, 0x80)      # 부제/캡션
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)     # 밝은 카드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=FONT, italic=False):
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return r


def _rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.5)
    return sp


def _oval(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.5)
    return sp


def _line(slide, x1, y1, x2, y2, color, weight=1.25):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(weight)
    cn.shadow.inherit = False
    return cn


def _topbar(slide, label=""):
    """본문 슬라이드 상단 액센트 바 + 모듈 라벨."""
    _rect(slide, 0, 0, 13.333, 0.16, fill=CORAL)
    if label:
        _, tf = _box(slide, 11.2, 0.28, 2.0, 0.4)
        _set(tf.paragraphs[0], label, 11, GRAY, bold=True, align=PP_ALIGN.RIGHT)


def _pagenum(slide, n):
    tb, tf = _box(slide, 12.4, 6.95, 0.8, 0.4)
    _set(tf.paragraphs[0], str(n), 11, GRAY, align=PP_ALIGN.RIGHT)


_counter = {"n": 0}


def new(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    _bg(s, bg)
    if _counter["n"] > 0:
        _pagenum(s, _counter["n"])
    _counter["n"] += 1
    return s


# ---------- 슬라이드 템플릿들 ----------

def cover(title, sub, foot):
    s = new(NAVY)
    _counter["n"] -= 1  # 표지엔 페이지번호 X (위에서 이미 증가했으므로 보정)
    # 위 보정이 어색하니 단순화: 표지는 직접 만든다
    return s


def slide_cover():
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    # 좌측 액센트 블록 + 코랄 점
    _rect(s, 0, 0, 0.35, 7.5, fill=CORAL)
    _oval(s, 11.7, 1.1, 1.0, 1.0, fill=CORAL)
    _oval(s, 12.2, 5.6, 0.6, 0.6, fill=GOLD)
    _rect(s, 1.0, 3.05, 13.333, 0.06, fill=CORAL)
    _, tf = _box(s, 1.0, 0.95, 6.0, 0.5)
    _set(tf.paragraphs[0], "AI 직원 만들기 · 마스터클래스", 15, GOLD, bold=True)
    _, tf = _box(s, 1.0, 1.9, 11.3, 1.3)
    _set(tf.paragraphs[0], "AI와 일하는 법", 56, WHITE, bold=True)
    _, tf = _box(s, 1.0, 3.25, 11.3, 1.0)
    _set(tf.paragraphs[0], "왜 AI는 강의를 할 수 없는가", 30, GOLD)
    _, tf = _box(s, 1.0, 5.8, 11.3, 0.6)
    _set(tf.paragraphs[0], "비개발자 셀러를 위한 3시간 · 코딩 없이 AI 직원 채용하기", 16, GRAY)


def slide_section(num, title, sub=""):
    s = new(NAVY)
    # 거대한 워터마크 숫자 (배경)
    _, tf = _box(s, 7.3, 1.4, 6.0, 5.2)
    tf.word_wrap = False
    _set(tf.paragraphs[0], str(num), 320, RGBColor(0x22, 0x2C, 0x44), bold=True, align=PP_ALIGN.RIGHT)
    _rect(s, 1.0, 2.35, 1.6, 0.06, fill=CORAL)
    _, tf = _box(s, 1.0, 2.5, 5.0, 1.4)
    _set(tf.paragraphs[0], f"MODULE {num}", 22, CORAL, bold=True)
    _, tf = _box(s, 1.0, 3.2, 9.5, 1.9)
    first = True
    for line in title.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        _set(p, line, 40, WHITE, bold=True)
        first = False
    if sub:
        _, tf = _box(s, 1.0, 5.3, 10.0, 1.0)
        _set(tf.paragraphs[0], sub, 19, GRAY)


def slide_statement(big, small="", accent=NAVY, bg=WHITE):
    """한 줄 강력 메시지."""
    s = new(bg)
    _, tf = _box(s, 1.0, 2.6, 11.3, 2.0)
    tf.word_wrap = True
    first = True
    for line in big.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        _set(p, line, 38, accent, bold=True, align=PP_ALIGN.CENTER)
        first = False
    if small:
        _, tf = _box(s, 1.5, 5.0, 10.3, 1.2)
        _set(tf.paragraphs[0], small, 19, GRAY, align=PP_ALIGN.CENTER)


def slide_title_bullets(title, bullets, kicker=""):
    """제목 + 불릿. bullets: list[(text, level, color)] 또는 list[str]."""
    s = new(WHITE)
    _topbar(s)
    if kicker:
        _, tf = _box(s, 1.0, 0.6, 11.3, 0.5)
        _set(tf.paragraphs[0], kicker, 15, CORAL, bold=True)
    _, tf = _box(s, 1.0, 1.05, 11.3, 1.1)
    _set(tf.paragraphs[0], title, 32, NAVY, bold=True)
    _rect(s, 1.02, 2.05, 0.9, 0.05, fill=GOLD)

    _, tf = _box(s, 1.0, 2.45, 11.3, 4.2)
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            text, level, color = b
        else:
            text, level, color = b, 0, INK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        size = 22 if level == 0 else 18
        bullet = "•  " if level == 0 else "–  "
        _set(p, bullet + text, size, color, bold=(level == 0 and color != INK))
        p.space_after = Pt(10)


def slide_compare(title, left_head, left_items, right_head, right_items):
    """좌(X) / 우(O) 2단 비교."""
    s = new(WHITE)
    _topbar(s)
    _, tf = _box(s, 1.0, 0.7, 11.3, 1.0)
    _set(tf.paragraphs[0], title, 30, NAVY, bold=True)

    # 색 띠로 좌우 성격 강조
    _rect(s, 0.9, 1.95, 5.5, 0.12, fill=CORAL)
    _rect(s, 6.9, 1.95, 5.5, 0.12, fill=TEAL)

    # 왼쪽 (부정)
    _rect(s, 0.9, 2.0, 5.5, 4.3, fill=RGBColor(0xFF, 0xEE, 0xEC))
    _, tf = _box(s, 1.2, 2.25, 5.0, 0.8)
    _set(tf.paragraphs[0], left_head, 22, CORAL, bold=True)
    _, tf = _box(s, 1.2, 3.1, 5.0, 3.0)
    first = True
    for it in left_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set(p, it, 18, INK)
        p.space_after = Pt(12)

    # 오른쪽 (긍정)
    _rect(s, 6.9, 2.0, 5.5, 4.3, fill=RGBColor(0xE9, 0xF7, 0xF2))
    _, tf = _box(s, 7.2, 2.25, 5.0, 0.8)
    _set(tf.paragraphs[0], right_head, 22, TEAL, bold=True)
    _, tf = _box(s, 7.2, 3.1, 5.0, 3.0)
    first = True
    for it in right_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set(p, it, 18, INK)
        p.space_after = Pt(12)


def slide_anatomy_diagram(title, parts):
    """AI 직원 해부도 — 가운데 직원 + 주변 7개 구성요소 다이어그램.
    parts: list[(라벨_한글, 실제용어, 색)] (7개)."""
    s = new(WHITE)
    _topbar(s)
    _, tf = _box(s, 1.0, 0.55, 11.3, 1.0)
    _set(tf.paragraphs[0], title, 30, NAVY, bold=True)

    # 중심 좌표 (제목 아래 공간의 중앙)
    cx, cy = 6.66, 4.4
    rx, ry = 4.45, 1.95   # 위성 배치 타원 반지름
    bw, bh = 2.5, 0.85    # 위성 박스 크기
    ocw, och = 2.3, 1.5   # 중심 타원 크기
    n = len(parts)

    # 1) 먼저 연결선(중심 → 각 위성). 박스보다 아래 레이어가 되도록 먼저 그림
    centers = []
    for i in range(n):
        ang = math.radians(-90 + i * (360.0 / n))
        bx_c = cx + rx * math.cos(ang)
        by_c = cy + ry * math.sin(ang)
        centers.append((bx_c, by_c))
        _line(s, cx, cy, bx_c, by_c, RGBColor(0xCF, 0xD6, 0xE0), weight=1.5)

    # 2) 중심 타원 (AI 직원)
    _oval(s, cx - ocw / 2, cy - och / 2, ocw, och, fill=NAVY)
    _, tf = _box(s, cx - ocw / 2, cy - 0.45, ocw, och)
    tf.anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], "AI 직원", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    p2 = tf.add_paragraph()
    _set(p2, "(내가 뽑은 직원)", 12, GOLD, align=PP_ALIGN.CENTER)

    # 3) 위성 박스들
    for (bx_c, by_c), (label, term, color) in zip(centers, parts):
        x, y = bx_c - bw / 2, by_c - bh / 2
        _rect(s, x, y, bw, bh, fill=LIGHT, line=color)
        _, tf = _box(s, x + 0.1, y + 0.06, bw - 0.2, bh)
        _set(tf.paragraphs[0], label, 16, color, bold=True, align=PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        _set(p2, term, 12, GRAY, align=PP_ALIGN.CENTER)


def slide_anatomy(title, rows):
    """AI 직원 해부도 / 번역표. rows: list[(left, right, note)]."""
    s = new(WHITE)
    _topbar(s)
    _, tf = _box(s, 1.0, 0.6, 11.3, 1.0)
    _set(tf.paragraphs[0], title, 30, NAVY, bold=True)
    _rect(s, 1.02, 1.6, 0.9, 0.05, fill=GOLD)

    y = 2.05
    rowh = 0.62
    # 헤더
    _, tf = _box(s, 1.1, y, 4.2, rowh)
    _set(tf.paragraphs[0], "쉬운 말 (강의 용어)", 15, GRAY, bold=True)
    _, tf = _box(s, 5.6, y, 3.2, rowh)
    _set(tf.paragraphs[0], "실제 용어", 15, GRAY, bold=True)
    _, tf = _box(s, 9.0, y, 3.6, rowh)
    _set(tf.paragraphs[0], "한 마디로", 15, GRAY, bold=True)
    y += rowh
    for i, (l, r, note) in enumerate(rows):
        if i % 2 == 0:
            _rect(s, 1.0, y, 11.4, rowh, fill=LIGHT)
        _, tf = _box(s, 1.1, y + 0.04, 4.4, rowh)
        _set(tf.paragraphs[0], l, 17, NAVY, bold=True)
        _, tf = _box(s, 5.6, y + 0.06, 3.2, rowh)
        _set(tf.paragraphs[0], r, 15, GRAY)
        _, tf = _box(s, 9.0, y + 0.06, 3.5, rowh)
        _set(tf.paragraphs[0], note, 15, INK)
        y += rowh


def slide_steps(title, steps):
    """가로 STEP 흐름."""
    s = new(WHITE)
    _topbar(s)
    _, tf = _box(s, 1.0, 0.7, 11.3, 1.0)
    _set(tf.paragraphs[0], title, 30, NAVY, bold=True)
    n = len(steps)
    cardw = 11.0 / n
    x = 1.15
    for i, (head, body) in enumerate(steps):
        _rect(s, x, 3.0, cardw - 0.25, 2.7, fill=LIGHT)
        # 번호 원형 배지
        cd = 0.66
        _oval(s, x + (cardw - 0.25) / 2 - cd / 2, 2.66, cd, cd, fill=CORAL)
        _, tf = _box(s, x + (cardw - 0.25) / 2 - cd / 2, 2.72, cd, cd)
        tf.anchor = MSO_ANCHOR.MIDDLE
        _set(tf.paragraphs[0], str(i + 1), 20, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _, tf = _box(s, x + 0.15, 3.55, cardw - 0.5, 0.9)
        _set(tf.paragraphs[0], head, 19, NAVY, bold=True, align=PP_ALIGN.CENTER)
        _, tf = _box(s, x + 0.15, 4.45, cardw - 0.5, 1.1)
        _set(tf.paragraphs[0], body, 14, GRAY, align=PP_ALIGN.CENTER)
        if i < n - 1:
            _, tf = _box(s, x + cardw - 0.34, 3.85, 0.5, 0.8)
            _set(tf.paragraphs[0], "›", 30, GOLD, bold=True)
        x += cardw


def slide_quote(text, who=""):
    s = new(NAVY)
    _, tf = _box(s, 1.3, 2.4, 10.7, 2.6)
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        _set(p, line, 30, WHITE, align=PP_ALIGN.CENTER, italic=True)
        first = False
    if who:
        _, tf = _box(s, 1.3, 5.2, 10.7, 0.6)
        _set(tf.paragraphs[0], "— " + who, 18, GOLD, align=PP_ALIGN.CENTER)


# ===================== 슬라이드 빌드 =====================

# --- 오프닝 ---
slide_cover()
slide_statement("이 강의는 필요 없습니다.", "정말로요. 「AI 직원 만들기」를 설치하면, 강의 없이도 작동합니다.", accent=NAVY)
slide_statement("이 슬라이드도\nAI가 만들었습니다.", "그런데 왜 사람이 강의를 할까요? — 오늘 그 답을 찾습니다.", accent=CORAL)
slide_title_bullets(
    "오늘의 약속",
    [
        ("코딩은 한 줄도 안 가르칩니다.", 0, INK),
        ("용어를 외우게 하지 않습니다.", 0, INK),
        ("대신, AI와 함께 '판단하는 법'을 가르칩니다.", 0, TEAL),
        ("끝나면, 당신은 AI 직원을 직접 채용할 수 있습니다.", 0, NAVY),
    ],
    kicker="이 강의가 하는 일",
)
slide_title_bullets(
    "오늘의 순서",
    [
        ("MODULE 1 — 왜 AI는 강의를 할 수 없는가  (사고방식)", 0, NAVY),
        ("MODULE 2 — AI 직원 해부도  (꼭 필요한 기본 용어)", 0, NAVY),
        ("MODULE 3 — 실전: AI 직원 채용하기  (스킬 사용)", 0, NAVY),
        ("그리고 — 운전면허 학원 이야기", 1, GRAY),
    ],
    kicker="목차",
)

# --- MODULE 1 ---
slide_section(1, "왜 AI는 강의를\n할 수 없는가", "AI가 다 해주는 시대, 그래도 사람이 배워야 하는 것")
slide_statement('"AI한테 만들어달라고\n하면 되는 거 아니에요?"', "맞습니다. 그게 함정입니다.", accent=NAVY)
slide_compare(
    "말 한마디가 결과를 바꿉니다",
    "❌ 명령형 — \"~해줘\"",
    ["\"상품등록 AI 만들어줘\"", "\"쿠팡 자동화 만들어줘\"", "\"광고 분석 툴 만들어줘\""],
    "⭕ 목표형 — \"~하고 싶다\"",
    ["\"상품등록 시간을 줄이고 싶다\"", "\"광고비 새는 데를 알고 싶다\"", "\"반복 문의 답변을 줄이고 싶다\""],
)
slide_title_bullets(
    "명령형의 함정",
    [
        ("\"만들어줘\"라고 하면 AI는 바로 만들기 시작합니다.", 0, INK),
        ("당신의 상황을 묻지 않습니다.", 1, GRAY),
        ("그래서 엉뚱한 결과가 나옵니다.", 0, CORAL),
        ("결국 처음부터 다시. 시간이 더 듭니다.", 1, GRAY),
    ],
    kicker="MODULE 1",
)
slide_title_bullets(
    "목표형의 힘",
    [
        ("\"~하고 싶다\"라고 하면 AI는 먼저 질문합니다.", 0, INK),
        ("\"지금은 어떻게 하세요?\" \"몇 명이 쓰나요?\"", 1, GRAY),
        ("그래서 당신 상황에 맞게 설계됩니다.", 0, TEAL),
        ("당신은 고르기만 하면 됩니다.", 1, GRAY),
    ],
    kicker="MODULE 1",
)
slide_statement("그럼 AI가 다 하네?\n— 아니요.", "AI가 구조적으로 '못 하는' 3가지가 있습니다.", accent=NAVY)
slide_title_bullets(
    "AI가 못 하는 것 ①  맥락",
    [
        ("AI는 '당신 회사'를 모릅니다.", 0, CORAL),
        ("직원 3명짜리와 30명짜리는 답이 다릅니다.", 0, INK),
        ("당신이 말해주지 않으면, AI는 '평균값'을 뱉습니다.", 0, INK),
        ("→ 그래서 AI는 당신에게 물어봐야 합니다.", 1, GRAY),
    ],
    kicker="AI가 못 하는 3가지",
)
slide_title_bullets(
    "AI가 못 하는 것 ②  판단",
    [
        ("AI는 A안과 B안을 줍니다.", 0, INK),
        ("하지만 '당신에게 뭐가 맞는지'는 정하지 못합니다.", 0, CORAL),
        ("고르는 것은 언제나 당신의 몫입니다.", 0, INK),
        ("→ 그래서 '고르는 법'을 배워야 합니다.", 1, GRAY),
    ],
    kicker="AI가 못 하는 3가지",
)
slide_title_bullets(
    "AI가 못 하는 것 ③  책임",
    [
        ("AI 결과가 틀려도, AI는 손해 보지 않습니다.", 0, INK),
        ("잘못 올라간 상품, 잘못 쓴 광고비 — 당신 몫입니다.", 0, CORAL),
        ("그래서 '언제 믿고 언제 멈출지'를 알아야 합니다.", 0, INK),
        ("→ 이게 오늘 강의의 마지막 주제입니다.", 1, GRAY),
    ],
    kicker="AI가 못 하는 3가지",
)
slide_statement("강의 = 맥락·판단·책임을\n당신에게 옮기는 일", "기능 설명이 아닙니다. AI가 못 하는 영역만 정확히 가르칩니다.", accent=TEAL)
slide_statement("당신은 코더가 아닙니다.\n당신은 AI 사장님입니다.", "AI 직원을 한 명 뽑았다고 생각하세요.", accent=NAVY)
slide_title_bullets(
    "사장님이 하는 일",
    [
        ("① 채용 — 어떤 직원이 필요한지 정한다", 0, NAVY),
        ("② 지시 — 무엇을 원하는지 목표로 말한다", 0, NAVY),
        ("③ 검수 — 결과가 맞는지 확인한다", 0, NAVY),
        ("코딩은? 직원(AI)이 합니다.", 0, TEAL),
    ],
    kicker="AI 사장님의 3가지 일",
)
slide_statement("MODULE 1 정리", "AI에게 '해줘' 대신 '하고 싶다'. AI가 못 하는 건 맥락·판단·책임. 당신은 사장님이다.", accent=NAVY)

# --- MODULE 2 ---
slide_section(2, "AI 직원 해부도", "꼭 필요한 기본 용어 — 외우지 말고, 직원 한 명을 떠올리세요")
slide_title_bullets(
    "용어가 무서운 진짜 이유",
    [
        ("MCP, 프롬프트, 워크플로우… 다 영어라서 무섭습니다.", 0, INK),
        ("그런데 사실은 회사에 이미 다 있는 개념입니다.", 0, TEAL),
        ("직원 한 명을 뜯어보면, 용어가 저절로 이해됩니다.", 0, NAVY),
    ],
    kicker="MODULE 2",
)
slide_anatomy_diagram(
    "AI 직원은 이렇게 생겼습니다",
    [
        ("두뇌", "Model", NAVY),
        ("기억", "Context", TEAL),
        ("손발", "Tool", GOLD),
        ("외부 연결", "MCP", CORAL),
        ("업무지침", "Instruction", NAVY),
        ("업무절차", "Workflow", TEAL),
        ("전문 능력", "Skill", CORAL),
    ],
)
slide_title_bullets(
    "두뇌 = AI의 머리 (모델)",
    [
        ("생각하고 글을 쓰는 부분. 사람으로 치면 머리입니다.", 0, INK),
        ("더 좋은 두뇌일수록 똑똑하지만, 비용도 큽니다.", 1, GRAY),
        ("→ 사장님은 '얼마나 똑똑한 직원이 필요한지'만 정하면 됩니다.", 0, TEAL),
    ],
    kicker="해부도 ①",
)
slide_title_bullets(
    "기억 = 상황·맥락 (Context)",
    [
        ("지금까지 나눈 대화와 상황을 기억하는 부분.", 0, INK),
        ("기억이 없으면 매번 처음부터 설명해야 합니다.", 1, GRAY),
        ("→ 그래서 '내 상황'을 잘 알려주는 게 중요합니다.", 0, TEAL),
    ],
    kicker="해부도 ②",
)
slide_title_bullets(
    "손발 & 외부 연결 = 도구 / MCP",
    [
        ("도구(Tool) = AI가 실제로 누르는 손발. 예: 파일 저장.", 0, INK),
        ("외부 서비스 연결(MCP) = 구글시트·쇼핑몰에 손을 뻗는 것.", 0, INK),
        ("→ '어떤 서비스랑 연결할까?'만 정하면 됩니다.", 0, TEAL),
    ],
    kicker="해부도 ③④",
)
slide_title_bullets(
    "매뉴얼 & 순서 = 업무지침 / 업무절차",
    [
        ("업무지침(Instruction) = 직원이 지켜야 할 규칙. \"존댓말로 답해라\".", 0, INK),
        ("업무절차(Workflow) = 일하는 순서. \"먼저 확인 → 그다음 등록\".", 0, INK),
        ("→ 신입에게 매뉴얼과 순서를 주는 것과 똑같습니다.", 0, TEAL),
    ],
    kicker="해부도 ⑤⑥",
)
slide_title_bullets(
    "전문 능력 = 스킬 (Skill)",
    [
        ("특정 분야를 잘하게 해주는 '자격증' 같은 것.", 0, INK),
        ("오늘 쓰는 「AI 직원 만들기」도 하나의 스킬입니다.", 0, NAVY),
        ("→ 필요한 능력만 골라 직원에게 달아줍니다.", 0, TEAL),
    ],
    kicker="해부도 ⑦",
)
slide_anatomy(
    "용어 번역표 (한 장 총정리)",
    [
        ("AI 직원 / 비서", "Agent", "일을 대신하는 AI"),
        ("대화", "Prompt", "AI에게 하는 말"),
        ("업무지침", "Instruction", "지켜야 할 규칙"),
        ("업무절차", "Workflow", "일하는 순서"),
        ("외부 서비스 연결", "MCP", "다른 서비스 연결"),
        ("도구", "Tool", "AI의 손발"),
        ("전문 능력", "Skill", "분야별 자격증"),
    ],
)
slide_statement("용어를 외우지 마세요.\n직원 한 명을 떠올리세요.", "두뇌·기억·손발·매뉴얼·순서·자격증. 그게 전부입니다.", accent=NAVY)
slide_statement("MODULE 2 정리", "AI 직원 = 두뇌+기억+손발+매뉴얼+순서+자격증. 용어는 회사에 이미 있던 개념이다.", accent=NAVY)

# --- MODULE 3 ---
slide_section(3, "실전: AI 직원\n채용하기", "이제 직접 한 명 뽑아봅니다 — 「AI 직원 만들기」 스킬")
slide_steps(
    "전체 흐름은 이게 전부입니다",
    [
        ("설치", "한 번만 하면 됨"),
        ("실행", "\"무엇이 불편하세요?\""),
        ("답하기", "목표·상황만 말함"),
        ("고르기", "추천안 중 선택"),
        ("완성", "직원이 일 시작"),
    ],
)
slide_title_bullets(
    "STEP 1 — 설치 (한 번만)",
    [
        ("「AI 직원 만들기」를 한 번만 설치합니다.", 0, INK),
        ("Claude Code 또는 Codex 어디서든 됩니다.", 1, GRAY),
        ("설치 후엔 그냥 부르기만 하면 됩니다.", 0, TEAL),
    ],
    kicker="MODULE 3",
)
slide_title_bullets(
    "STEP 2 — 실행하면 코치가 먼저 묻습니다",
    [
        ("\"안녕하세요. 코딩은 몰라도 됩니다.\"", 0, NAVY),
        ("\"요즘 어떤 일이 제일 번거로우세요?\"", 0, NAVY),
        ("당신은 무엇을 만들지 몰라도 됩니다. 불편한 것만 말하세요.", 0, TEAL),
    ],
    kicker="MODULE 3",
)
slide_title_bullets(
    "목표 말하는 법 — 불편/욕구만",
    [
        ("\"상품등록이 너무 오래 걸려요\"", 0, INK),
        ("\"광고비가 어디서 새는지 모르겠어요\"", 0, INK),
        ("\"같은 고객 문의에 매번 답하는 게 지쳐요\"", 0, INK),
        ("이렇게만 말하면, 나머지는 코치가 질문합니다.", 0, TEAL),
    ],
    kicker="이렇게 말하세요",
)
slide_title_bullets(
    "AI의 질문에 답하는 법",
    [
        ("코치는 한 번에 하나씩 묻습니다.", 0, INK),
        ("\"지금은 어떻게 하세요?\" → 솔직하게 답하면 됩니다.", 1, GRAY),
        ("모르면 \"잘 모르겠어요\"도 정답입니다.", 0, TEAL),
        ("좋은 AI는 질문을 많이 하는 AI입니다.", 0, NAVY),
    ],
    kicker="MODULE 3",
)
slide_title_bullets(
    "판단 시뮬레이터 — A/B가 나오면?",
    [
        ("코치: \"추천 ① 반자동 — AI가 초안까지, 등록 버튼은 사장님이\"", 0, TEAL),
        ("코치: \"② 완전자동 — AI가 등록까지. 단, 되돌리기 번거로움\"", 0, INK),
        ("당신: 둘 중 하나만 고르면 됩니다.", 0, NAVY),
        ("정답은 없습니다. '내 상황에 맞는' 게 정답입니다.", 1, GRAY),
    ],
    kicker="강의의 핵심",
)
slide_title_bullets(
    "고르는 기준",
    [
        ("처음이라면 → 항상 '추천안(보통 ①)'으로.", 0, TEAL),
        ("규모가 커지면 → 더 자동화된 쪽으로 옮긴다.", 0, INK),
        ("불안하면 → 사람이 마지막을 확인하는 쪽으로.", 0, INK),
        ("이 기준만 있으면, 어떤 선택지가 나와도 고를 수 있습니다.", 0, NAVY),
    ],
    kicker="이게 평생 가는 능력",
)
slide_title_bullets(
    "결과 검수하는 법",
    [
        ("완성되면, 작은 것 하나로 먼저 테스트합니다.", 0, INK),
        ("\"상품 1개만 먼저 해보자\" — 전체에 돌리기 전에.", 1, GRAY),
        ("결과가 내 기대와 같은지 눈으로 확인합니다.", 0, TEAL),
    ],
    kicker="MODULE 3",
)
slide_statement("🛑 멈춤 버튼", "다들 '어떻게 시키나'만 배웁니다. 진짜 중요한 건 '언제 멈추나'입니다.", accent=CORAL)
slide_compare(
    "언제 믿고, 언제 직접 볼까",
    "🛑 멈추고 직접 확인",
    ["돈이 나가는 일 (결제·광고비)", "되돌리기 어려운 일 (삭제·발송)", "결과가 평소와 많이 다를 때", "AI가 확신 없이 얼버무릴 때"],
    "✅ 믿고 맡겨도 OK",
    ["초안·아이디어 만들기", "반복되는 단순 작업", "언제든 되돌릴 수 있는 일", "이미 여러 번 검증된 작업"],
)
slide_statement("MODULE 3 정리", "목표만 말하고, 추천안 고르고, 작게 테스트하고, 위험하면 멈춘다.", accent=NAVY)

# --- 클로징 ---
slide_section("∞", "다시, 왜 AI는\n강의를 할 수 없는가", "이제 답할 수 있습니다")
slide_title_bullets(
    "AI는 정보를 줍니다. 하지만…",
    [
        ("AI는 '당신이 무엇을 모르는지'를 모릅니다.", 0, INK),
        ("AI는 '당신 회사의 사정'을 모릅니다.  (맥락)", 0, INK),
        ("AI는 '당신 대신 책임'지지 않습니다.  (책임)", 0, INK),
        ("그래서 판단은 끝까지 사람의 몫입니다.  (판단)", 0, CORAL),
    ],
    kicker="오늘의 답",
)
slide_quote(
    "이 강의의 경쟁자는\nAI 강의가 아니라 운전면허 학원입니다.",
    "운전학원은 엔진 구조를 가르치지 않습니다",
)
slide_title_bullets(
    "운전면허 학원처럼",
    [
        ("엔진 구조(MCP가 뭔지)는 깊게 안 배웁니다.", 0, GRAY),
        ("브레이크 언제 밟나 = 언제 AI를 멈추나", 0, NAVY),
        ("차선 바꾸기 = 방식을 언제 바꾸나", 0, NAVY),
        ("위험 판단 = 언제 의심하나", 0, NAVY),
    ],
    kicker="우리가 가르친 것",
)
slide_statement("도구는 바뀝니다.\n판단하는 법은 남습니다.", "스킬은 1년 뒤 바뀔 수 있습니다. 오늘 배운 판단력은 그대로 갑니다.", accent=TEAL)
slide_title_bullets(
    "당신의 다음 한 걸음",
    [
        ("① 「AI 직원 만들기」를 설치한다.", 0, NAVY),
        ("② 가장 번거로운 일 하나를 떠올린다.", 0, NAVY),
        ("③ 코치에게 \"이게 불편해요\"라고 말한다.", 0, NAVY),
        ("그게 당신의 첫 AI 직원입니다.", 0, TEAL),
    ],
    kicker="지금 바로",
)
slide_statement("당신은 이제\nAI 사장님입니다.", "고맙습니다.", accent=NAVY)

prs.save("AI와_일하는_법.pptx")
print(f"saved: AI와_일하는_법.pptx  ({len(prs.slides)} slides)")
